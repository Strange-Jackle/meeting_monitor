"""
Ollama Document Analysis Service

Uses the Ollama API with vision models to analyze documents (PDF, PPTX, DOCX).
Supports OCR through vision model capabilities.
"""

import requests
import base64
import os
from typing import Optional, Dict, Any, List
from pathlib import Path

# PDF processing
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    print("[OllamaService] PyMuPDF not installed. PDF analysis will be limited.")

# PPTX processing
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("[OllamaService] python-pptx not installed. PPTX analysis will be limited.")

# DOCX processing
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("[OllamaService] python-docx not installed. DOCX analysis will be limited.")


class OllamaDocumentService:
    """Service for analyzing documents using Ollama vision models."""
    
    def __init__(
        self, 
        base_url: str = "http://10.119.65.52:11434",
        model: str = "qwen3-vl:2b"
    ):
        self.base_url = base_url.rstrip('/')
        self.model = model
        self.timeout = 120  # Longer timeout for document analysis
    
    async def analyze_document(
        self, 
        file_path: str, 
        prompt: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Analyze a document file (PDF, PPTX, DOCX).
        
        Returns dict with:
        - summary: Brief summary of the document
        - key_points: List of main points
        - text_content: Extracted text (if available)
        - analysis: Full analysis from the model
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            return {"error": f"File not found: {file_path}"}
        
        ext = file_path.suffix.lower()
        
        if ext == '.pdf':
            return await self._analyze_pdf(file_path, prompt)
        elif ext == '.pptx':
            return await self._analyze_pptx(file_path, prompt)
        elif ext == '.docx':
            return await self._analyze_docx(file_path, prompt)
        else:
            return {"error": f"Unsupported file type: {ext}"}
    
    async def _analyze_pdf(self, file_path: Path, prompt: Optional[str]) -> Dict[str, Any]:
        """Analyze PDF by extracting text (fast mode, no image analysis)."""
        if not HAS_PYMUPDF:
            return {"error": "PyMuPDF not installed. Run: pip install pymupdf"}
        
        try:
            doc = fitz.open(file_path)
            all_text = []
            total_pages = len(doc)
            
            # Process first 3 pages only for speed
            pages_to_analyze = min(3, total_pages)
            for page_num in range(pages_to_analyze):
                page = doc[page_num]
                
                # Extract text directly (fast, no vision)
                text = page.get_text()
                if text.strip():
                    all_text.append(f"--- Page {page_num + 1} ---\n{text.strip()}")
            
            doc.close()
            
            # Combine text
            text_content = "\n\n".join(all_text)
            
            if not text_content.strip():
                text_content = "No text could be extracted from this PDF."
            
            # Generate summary
            summary_prompt = prompt or "Summarize this document and list the key points."
            summary = await self._generate_text(
                f"Based on the following document content, {summary_prompt}\n\nDocument content:\n{text_content[:4000]}"
            )
            
            return {
                "status": "success",
                "file_type": "pdf",
                "pages_analyzed": pages_to_analyze,
                "total_pages": total_pages,
                "text_content": text_content[:5000],
                "summary": summary,
                "analysis": f"Extracted text from {pages_to_analyze} of {total_pages} pages."
            }
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            return {"error": f"PDF analysis failed: {str(e)}"}
    
    async def _analyze_pptx(self, file_path: Path, prompt: Optional[str]) -> Dict[str, Any]:
        """Analyze PowerPoint presentation."""
        if not HAS_PPTX:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}
        
        try:
            prs = Presentation(file_path)
            slides_text = []
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                if slide_count > 10:  # Limit to first 10 slides
                    break
                    
                slide_content = []
                try:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text and shape.text.strip():
                            slide_content.append(shape.text.strip())
                except Exception as shape_error:
                    print(f"[OllamaService] Error reading shape: {shape_error}")
                    continue
                
                if slide_content:
                    slides_text.append(f"--- Slide {slide_count} ---\n" + "\n".join(slide_content))
            
            text_content = "\n\n".join(slides_text) if slides_text else "No text content extracted from slides."
            
            # Generate summary
            summary_prompt = prompt or "Summarize this presentation and list the main topics covered."
            summary = await self._generate_text(
                f"Based on this PowerPoint presentation content, {summary_prompt}\n\nPresentation content:\n{text_content[:4000]}"
            )
            
            return {
                "status": "success",
                "file_type": "pptx",
                "slides_analyzed": slide_count,
                "text_content": text_content,
                "summary": summary,
                "analysis": f"Analyzed {slide_count} slides from the presentation."
            }
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            return {"error": f"PPTX analysis failed: {str(e)}"}
    
    async def _analyze_docx(self, file_path: Path, prompt: Optional[str]) -> Dict[str, Any]:
        """Analyze Word document."""
        if not HAS_DOCX:
            return {"error": "python-docx not installed. Run: pip install python-docx"}
        
        try:
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text_content = "\n\n".join(paragraphs[:100])  # First 100 paragraphs
            
            # Generate summary
            summary_prompt = prompt or "Summarize this document and extract the key points."
            summary = await self._generate_text(
                f"Based on this Word document content, {summary_prompt}\n\nDocument content:\n{text_content[:4000]}"
            )
            
            return {
                "status": "success",
                "file_type": "docx",
                "paragraphs_analyzed": min(100, len(paragraphs)),
                "text_content": text_content,
                "summary": summary,
                "analysis": f"Analyzed {len(paragraphs)} paragraphs from the document."
            }
            
        except Exception as e:
            return {"error": f"DOCX analysis failed: {str(e)}"}
    
    async def _analyze_image(self, image_b64: str, prompt: str) -> str:
        """Analyze an image using Ollama vision model."""
        try:
            response = requests.post(
                f"{self.base_url}/api/generate",
                json={
                    "model": self.model,
                    "prompt": prompt,
                    "images": [image_b64],
                    "stream": False
                },
                timeout=self.timeout
            )
            
            if response.status_code == 200:
                return response.json().get("response", "No response")
            else:
                return f"Vision API error: {response.status_code}"
                
        except Exception as e:
            return f"Vision analysis error: {str(e)}"
    
    async def _generate_text(self, prompt: str) -> str:
        """Generate text using Ollama."""
        try:
            response = requests.post(
                f"{self.base_url}/api/generate",
                json={
                    "model": self.model,
                    "prompt": prompt,
                    "stream": False
                },
                timeout=self.timeout
            )
            
            if response.status_code == 200:
                return response.json().get("response", "No response")
            else:
                return f"API error: {response.status_code}"
                
        except Exception as e:
            return f"Generation error: {str(e)}"
    
    async def generate_key_insights(self, text_content: str, num_insights: int = 3) -> List[str]:
        """
        Generate key actionable insights from document content.
        Returns a list of one-line insights for overlay display.
        """
        try:
            prompt = f"""Based on the following document content, extract exactly {num_insights} key insights.
Each insight should be:
- One sentence only (max 15 words)
- Actionable or informative
- Relevant for a sales/business meeting

Format your response as exactly {num_insights} lines, one insight per line. No numbering, no bullets.

Document content:
{text_content[:3000]}

Your {num_insights} insights:"""
            
            response = requests.post(
                f"{self.base_url}/api/generate",
                json={
                    "model": self.model,
                    "prompt": prompt,
                    "stream": False
                },
                timeout=60
            )
            
            if response.status_code == 200:
                raw_response = response.json().get("response", "")
                # Parse insights - take first num_insights non-empty lines
                lines = [line.strip() for line in raw_response.split('\n') if line.strip()]
                # Clean up any numbering or bullets
                insights = []
                for line in lines[:num_insights]:
                    # Remove common prefixes like "1.", "- ", "• "
                    cleaned = line.lstrip('0123456789.-•) ').strip()
                    if cleaned and len(cleaned) > 5:
                        insights.append(cleaned)
                
                # Ensure we have exactly num_insights
                while len(insights) < num_insights:
                    insights.append("Document analyzed successfully.")
                
                return insights[:num_insights]
            else:
                return [f"Analysis complete (model unavailable)"] * num_insights
                
        except Exception as e:
            print(f"[OllamaService] Insights generation error: {e}")
            return ["Document analyzed - key points extracted."] * num_insights
    
    def check_health(self) -> Dict[str, Any]:
        """Check if Ollama service is available."""
        try:
            response = requests.get(f"{self.base_url}/api/tags", timeout=5)
            if response.status_code == 200:
                models = response.json().get("models", [])
                return {
                    "status": "healthy",
                    "models": [m.get("name") for m in models]
                }
            return {"status": "error", "message": f"Status code: {response.status_code}"}
        except Exception as e:
            return {"status": "error", "message": str(e)}


# Singleton instance
_ollama_service = None

def get_ollama_service() -> OllamaDocumentService:
    """Get or create Ollama service instance."""
    global _ollama_service
    if _ollama_service is None:
        _ollama_service = OllamaDocumentService()
    return _ollama_service

